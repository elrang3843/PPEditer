"""
Slide renderer — converts python-pptx Slide objects to QPixmap using QPainter.

Font-size scaling:
    display_pt = font_pt * scale_y * EMU_PER_INCH / SCREEN_DPI
where scale_y = display_height_px / slide_height_emu
"""
from PyQt5.QtCore import Qt, QRect
from PyQt5.QtGui import (
    QAbstractTextDocumentLayout,
    QColor,
    QFont,
    QImage,
    QPainter,
    QPen,
    QPixmap,
    QTextDocument,
)

EMU_PER_INCH = 914400
EMU_PER_PT = 12700
SCREEN_DPI = 96
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# ── Public API ─────────────────────────────────────────────────────────────────

def render_slide(slide, width_px: int, height_px: int,
                 slide_width_emu: int, slide_height_emu: int) -> QPixmap:
    """Render *slide* into a QPixmap of size (*width_px* × *height_px*)."""
    pixmap = QPixmap(max(1, width_px), max(1, height_px))
    pixmap.fill(Qt.white)
    if slide is None:
        return pixmap

    painter = QPainter(pixmap)
    painter.setRenderHint(QPainter.Antialiasing)
    painter.setRenderHint(QPainter.TextAntialiasing)
    painter.setRenderHint(QPainter.SmoothPixmapTransform)

    sx = width_px / max(1, slide_width_emu)
    sy = height_px / max(1, slide_height_emu)

    _draw_background(painter, slide, width_px, height_px)
    for shape in slide.shapes:
        try:
            _draw_shape(painter, shape, sx, sy)
        except Exception:
            pass

    painter.end()
    return pixmap


# ── Background ─────────────────────────────────────────────────────────────────

def _draw_background(painter: QPainter, slide, w: int, h: int):
    try:
        fill = slide.background.fill
        if fill.type is not None:
            c = _rgb_color(fill.fore_color.rgb)
            if c:
                painter.fillRect(0, 0, w, h, c)
                return
    except Exception:
        pass
    painter.fillRect(0, 0, w, h, Qt.white)


# ── Shape dispatch ─────────────────────────────────────────────────────────────

def _shape_rect(shape, sx: float, sy: float) -> QRect:
    return QRect(
        int((shape.left or 0) * sx),
        int((shape.top or 0) * sy),
        int((shape.width or 0) * sx),
        int((shape.height or 0) * sy),
    )


def _draw_shape(painter: QPainter, shape, sx: float, sy: float):
    rect = _shape_rect(shape, sx, sy)
    if rect.width() <= 0 or rect.height() <= 0:
        return

    _draw_fill(painter, shape, rect)
    _draw_border(painter, shape, rect)

    # Picture
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _draw_picture(painter, shape, rect)
            return
    except Exception:
        pass

    if shape.has_text_frame:
        _draw_text(painter, shape, rect, sy)


# ── Fill / border ──────────────────────────────────────────────────────────────

def _draw_fill(painter: QPainter, shape, rect: QRect):
    try:
        fill = shape.fill
        if fill.type is not None:
            c = _rgb_color(fill.fore_color.rgb)
            if c:
                painter.fillRect(rect, c)
    except Exception:
        pass


def _draw_border(painter: QPainter, shape, rect: QRect):
    try:
        line = shape.line
        if line.color.type is not None:
            c = _rgb_color(line.color.rgb)
            if c:
                w = max(1, int((line.width or EMU_PER_PT) / EMU_PER_PT))
                pen = QPen(c, w)
                old = painter.pen()
                painter.setPen(pen)
                painter.drawRect(rect)
                painter.setPen(old)
    except Exception:
        pass


# ── Picture ────────────────────────────────────────────────────────────────────

def _draw_picture(painter: QPainter, shape, rect: QRect):
    try:
        img = QImage()
        if img.loadFromData(shape.image.blob):
            painter.drawImage(rect, img)
    except Exception:
        pass


# ── Text ───────────────────────────────────────────────────────────────────────

def _draw_text(painter: QPainter, shape, rect: QRect, sy: float):
    try:
        tf = shape.text_frame
        html = _build_html(tf, sy)

        doc = QTextDocument()
        doc.setHtml(html)
        doc.setTextWidth(rect.width())

        doc_h = doc.size().height()

        y_offset = 0.0
        try:
            from pptx.enum.text import MSO_ANCHOR
            anchor = tf.vertical_anchor
            if anchor == MSO_ANCHOR.MIDDLE:
                y_offset = max(0.0, (rect.height() - doc_h) / 2)
            elif anchor == MSO_ANCHOR.BOTTOM:
                y_offset = max(0.0, rect.height() - doc_h)
        except Exception:
            pass

        painter.save()
        painter.setClipRect(rect)
        painter.translate(rect.left(), rect.top() + y_offset)
        ctx = QAbstractTextDocumentLayout.PaintContext()
        doc.documentLayout().draw(painter, ctx)
        painter.restore()

    except Exception:
        # Plain-text fallback
        try:
            text = shape.text_frame.text
            if text:
                painter.save()
                painter.setPen(Qt.black)
                painter.setFont(QFont("맑은 고딕", max(6, int(14 * sy * EMU_PER_INCH / SCREEN_DPI))))
                painter.drawText(rect, Qt.AlignLeft | Qt.AlignTop | Qt.TextWordWrap, text)
                painter.restore()
        except Exception:
            pass


# ── HTML builder ───────────────────────────────────────────────────────────────

def _build_html(tf, sy: float) -> str:
    parts = ['<html><body style="margin:0;padding:0;">']
    for para in tf.paragraphs:
        align = _para_align(para)
        parts.append(f'<p style="text-align:{align};margin:0;padding:0;">')
        if not para.runs:
            # Empty paragraph — reserve vertical space using default size
            sz = _emu_to_display_pt(18 * EMU_PER_PT, sy)
            parts.append(f'<span style="font-size:{sz:.1f}pt">&nbsp;</span>')
        for run in para.runs:
            style = _run_css(run, sy)
            text = (run.text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            parts.append(f'<span style="{style}">{text}</span>')
        parts.append("</p>")
    parts.append("</body></html>")
    return "".join(parts)


def _para_align(para) -> str:
    try:
        from pptx.enum.text import PP_ALIGN
        return {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
        }.get(para.alignment, "left")
    except Exception:
        return "left"


def _run_css(run, sy: float) -> str:
    parts = []
    try:
        f = run.font
        parts.append(f'font-family:"{f.name or "맑은 고딕"}"')
        size_emu = f.size or (18 * EMU_PER_PT)
        pt = _emu_to_display_pt(size_emu, sy)
        parts.append(f"font-size:{max(1.0, pt):.1f}pt")
        if f.bold:
            parts.append("font-weight:bold")
        if f.italic:
            parts.append("font-style:italic")
        if f.underline:
            parts.append("text-decoration:underline")
        try:
            if f.color.type is not None:
                c = _rgb_color(f.color.rgb)
                if c:
                    parts.append(f"color:rgb({c.red()},{c.green()},{c.blue()})")
                else:
                    parts.append("color:black")
            else:
                parts.append("color:black")
        except Exception:
            parts.append("color:black")
    except Exception:
        return 'font-family:"맑은 고딕";font-size:14pt;color:black'
    return ";".join(parts)


# ── Helpers ────────────────────────────────────────────────────────────────────

def _emu_to_display_pt(size_emu: int, sy: float) -> float:
    """Convert EMU font size to display pt, accounting for current zoom."""
    return (size_emu / EMU_PER_PT) * sy * EMU_PER_INCH / SCREEN_DPI


def _rgb_color(rgb) -> QColor | None:
    try:
        return QColor(rgb.red, rgb.green, rgb.blue)
    except Exception:
        return None

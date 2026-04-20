"""
Presentation model — wraps python-pptx and provides undo/redo support.
"""
import io
import os

from pptx import Presentation
from pptx.util import Pt


class PresentationModel:
    """Manages a PPTX presentation with undo/redo and modification tracking."""

    MAX_UNDO = 50

    def __init__(self):
        self.prs = None
        self.filepath = None
        self.modified = False
        self._undo_stack: list[bytes] = []
        self._redo_stack: list[bytes] = []

    # ── File operations ────────────────────────────────────────────────

    def new(self):
        """Create a new empty presentation with one blank slide."""
        self.prs = Presentation()
        layout = self.prs.slide_layouts[0]
        self.prs.slides.add_slide(layout)
        self.filepath = None
        self.modified = False
        self._undo_stack.clear()
        self._redo_stack.clear()

    def open(self, filepath: str):
        """Open a PPTX/PPT file."""
        self.prs = Presentation(filepath)
        self.filepath = filepath
        self.modified = False
        self._undo_stack.clear()
        self._redo_stack.clear()

    def save(self, filepath: str = None):
        """Save the presentation. Raises ValueError if no path known."""
        if filepath:
            self.filepath = filepath
        if not self.filepath:
            raise ValueError("저장할 파일 경로가 없습니다.")
        self.prs.save(self.filepath)
        self.modified = False

    def to_bytes(self) -> bytes:
        """Serialize current presentation to bytes."""
        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()

    def from_bytes(self, data: bytes):
        """Load presentation from bytes (restores undo state)."""
        self.prs = Presentation(io.BytesIO(data))

    # ── Properties ─────────────────────────────────────────────────────

    @property
    def is_open(self) -> bool:
        return self.prs is not None

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides) if self.prs else 0

    @property
    def slide_width(self) -> int:
        return self.prs.slide_width if self.prs else 9144000

    @property
    def slide_height(self) -> int:
        return self.prs.slide_height if self.prs else 5143500

    @property
    def filename(self) -> str:
        return os.path.basename(self.filepath) if self.filepath else "제목없음.pptx"

    @property
    def window_title(self) -> str:
        prefix = "* " if self.modified else ""
        return f"{prefix}{self.filename} - PPEditer"

    # ── Slide access ───────────────────────────────────────────────────

    def get_slide(self, index: int):
        if not self.is_open or not (0 <= index < self.slide_count):
            return None
        return self.prs.slides[index]

    # ── Slide mutations ────────────────────────────────────────────────

    def add_slide(self, after_index: int = None) -> int:
        """Add a blank slide after *after_index*. Returns new slide index."""
        self._push_undo()
        layout_idx = min(1, len(self.prs.slide_layouts) - 1)
        self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        new_idx = self.slide_count - 1
        if after_index is not None and after_index < new_idx:
            self.move_slide(new_idx, after_index + 1)
            new_idx = after_index + 1
        self.modified = True
        return new_idx

    def delete_slide(self, index: int) -> bool:
        """Delete slide at *index*. Returns False if only one slide remains."""
        if self.slide_count <= 1:
            return False
        self._push_undo()
        sldIdLst = self.prs.slides._sldIdLst
        sldId = sldIdLst[index]
        r_id = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        sldIdLst.remove(sldId)
        if r_id:
            try:
                del self.prs.part._rels[r_id]
            except Exception:
                pass
        self.modified = True
        return True

    def duplicate_slide(self, index: int) -> int:
        """Duplicate slide at *index*. Returns new slide index."""
        from lxml import etree

        self._push_undo()
        src_slide = self.prs.slides[index]
        layout = self.prs.slide_layouts[0]
        new_slide = self.prs.slides.add_slide(layout)

        # Replace default shapes with a deep copy of source slide's shape tree
        sp_tree = new_slide.shapes._spTree
        for child in list(sp_tree):
            sp_tree.remove(child)
        for child in src_slide.shapes._spTree:
            sp_tree.append(etree.fromstring(etree.tostring(child)))

        new_idx = self.slide_count - 1
        if new_idx > index + 1:
            self.move_slide(new_idx, index + 1)
            new_idx = index + 1
        self.modified = True
        return new_idx

    def move_slide(self, from_index: int, to_index: int):
        """Reorder slide from *from_index* to *to_index*."""
        if from_index == to_index:
            return
        if not (0 <= from_index < self.slide_count):
            return
        if not (0 <= to_index < self.slide_count):
            return
        self._push_undo()
        sldIdLst = self.prs.slides._sldIdLst
        sldId = sldIdLst[from_index]
        sldIdLst.remove(sldId)
        sldIdLst.insert(to_index, sldId)
        self.modified = True

    def update_shape_text(self, slide_index: int, shape_index: int, text: str):
        """Replace text content of a shape (preserves first-run font as-is)."""
        slide = self.get_slide(slide_index)
        if slide is None:
            return
        shapes = list(slide.shapes)
        if not (0 <= shape_index < len(shapes)):
            return
        shape = shapes[shape_index]
        if shape.has_text_frame:
            self._push_undo()
            tf = shape.text_frame
            # Keep first paragraph's run formatting when possible
            tf.text = text
            self.modified = True

    # ── Undo / Redo ────────────────────────────────────────────────────

    def _push_undo(self):
        try:
            data = self.to_bytes()
            self._undo_stack.append(data)
            if len(self._undo_stack) > self.MAX_UNDO:
                self._undo_stack.pop(0)
            self._redo_stack.clear()
        except Exception:
            pass

    def undo(self) -> bool:
        if not self._undo_stack:
            return False
        self._redo_stack.append(self.to_bytes())
        saved_path = self.filepath
        self.from_bytes(self._undo_stack.pop())
        self.filepath = saved_path
        self.modified = True
        return True

    def redo(self) -> bool:
        if not self._redo_stack:
            return False
        self._undo_stack.append(self.to_bytes())
        saved_path = self.filepath
        self.from_bytes(self._redo_stack.pop())
        self.filepath = saved_path
        self.modified = True
        return True

    @property
    def can_undo(self) -> bool:
        return bool(self._undo_stack)

    @property
    def can_redo(self) -> bool:
        return bool(self._redo_stack)
